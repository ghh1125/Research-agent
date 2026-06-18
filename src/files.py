from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class ParsedFile:
    path: str
    kind: str
    text: str
    sheets: dict[str, list[list[str]]] = field(default_factory=dict)
    error: str | None = None


def parse_file(path: str | Path) -> ParsedFile:
    p = Path(path)
    if not p.exists():
        return ParsedFile(path=str(p), kind="missing", text="", error=f"file not found: {p}")
    suffix = p.suffix.lower()
    try:
        if suffix == ".pdf":
            return ParsedFile(path=str(p), kind="pdf", text=_parse_pdf(p))
        if suffix == ".docx":
            return ParsedFile(path=str(p), kind="docx", text=_parse_docx(p))
        if suffix in (".pptx", ".ppt"):
            return ParsedFile(path=str(p), kind="pptx", text=_parse_pptx(p))
        if suffix in (".xlsx", ".xls"):
            sheets = _parse_xlsx(p)
            flattened = "\n\n".join(f"[sheet:{name}]\n" + "\n".join(",".join(row) for row in rows) for name, rows in sheets.items())
            return ParsedFile(path=str(p), kind="xlsx", text=flattened, sheets=sheets)
        if suffix in (".txt", ".md"):
            return ParsedFile(path=str(p), kind="text", text=p.read_text(encoding="utf-8", errors="ignore"))
        return ParsedFile(path=str(p), kind="unsupported", text="", error=f"unsupported file type: {suffix}")
    except Exception as exc:
        return ParsedFile(path=str(p), kind=suffix.lstrip("."), text="", error=f"{type(exc).__name__}: {exc}")


def parse_files(paths: list[str | Path]) -> list[ParsedFile]:
    return [parse_file(path) for path in paths]


def _parse_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _parse_docx(path: Path) -> str:
    from docx import Document

    document = Document(str(path))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _parse_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"[slide {index}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def _parse_xlsx(path: Path) -> dict[str, list[list[str]]]:
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), data_only=True)
    sheets: dict[str, list[list[str]]] = {}
    for name in workbook.sheetnames:
        sheet = workbook[name]
        rows: list[list[str]] = []
        for row in sheet.iter_rows(values_only=True):
            rows.append(["" if cell is None else str(cell) for cell in row])
        sheets[name] = rows
    return sheets


def truncate(text: str, max_chars: int = 6000) -> str:
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + f"\n...[截断，原文 {len(text)} 字符]"


def save_uploaded_bytes(items: list[tuple[str, bytes]], work_dir: str | Path, subdir: str) -> list[str]:
    """Write (filename, content) pairs (e.g. from a web upload widget) under work_dir/subdir and return their paths."""

    if not items:
        return []
    target_dir = Path(work_dir) / subdir
    target_dir.mkdir(parents=True, exist_ok=True)
    paths: list[str] = []
    for filename, content in items:
        path = target_dir / filename
        path.write_bytes(content)
        paths.append(str(path))
    return paths
